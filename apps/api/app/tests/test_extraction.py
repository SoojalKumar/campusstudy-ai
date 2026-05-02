from io import BytesIO
from types import SimpleNamespace

from docx import Document
from pypdf import PdfWriter
from pptx import Presentation

from app.models.enums import MaterialKind
from app.providers.base import TranscriptionSegment
from app.services.extraction import MaterialExtractor, detect_material_kind


def build_docx_bytes() -> bytes:
    document = Document()
    document.add_paragraph("Sorting stability matters when equivalent keys must preserve order.")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Hash tables"
    slide.placeholders[1].text = "Average lookup is O(1) with a good hash distribution."
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def build_pdf_bytes() -> bytes:
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def test_extractor_reads_text_docx_and_pptx_sources():
    extractor = MaterialExtractor()

    text_material = SimpleNamespace(file_name="lecture.txt", source_kind=MaterialKind.DOCUMENT)
    text_content = extractor.extract(text_material, b"Queues support FIFO ordering.")
    assert text_content.full_text == "Queues support FIFO ordering."

    docx_material = SimpleNamespace(file_name="lecture.docx", source_kind=MaterialKind.DOCUMENT)
    docx_content = extractor.extract(docx_material, build_docx_bytes())
    assert "Sorting stability matters" in docx_content.full_text

    pptx_material = SimpleNamespace(file_name="slides.pptx", source_kind=MaterialKind.DOCUMENT)
    pptx_content = extractor.extract(pptx_material, build_pptx_bytes())
    assert any(section.slide_number == 1 for section in pptx_content.sections)
    assert "Hash tables" in pptx_content.full_text


def test_extractor_handles_pdf_and_audio_transcript_sources():
    extractor = MaterialExtractor()
    extractor.stt_provider = SimpleNamespace(
        transcribe=lambda _: [
            TranscriptionSegment(start_second=0, end_second=42, text="Lecture intro on regression assumptions."),
            TranscriptionSegment(start_second=43, end_second=85, text="Bias variance tradeoff and validation checks."),
        ]
    )

    pdf_material = SimpleNamespace(file_name="reader.pdf", source_kind=MaterialKind.DOCUMENT)
    pdf_content = extractor.extract(pdf_material, build_pdf_bytes())
    assert len(pdf_content.sections) == 1
    assert pdf_content.sections[0].page_number == 1

    audio_material = SimpleNamespace(file_name="lecture.mp3", source_kind=MaterialKind.AUDIO)
    transcript_content = extractor.extract(audio_material, b"fake-audio")
    assert len(transcript_content.transcript_segments) == 2
    assert transcript_content.sections[0].start_second == 0


def test_detect_material_kind_uses_extension_and_mime_type():
    assert detect_material_kind("lecture.mp3", "audio/mpeg") == MaterialKind.AUDIO
    assert detect_material_kind("recording.bin", "video/mp4") == MaterialKind.VIDEO
    assert detect_material_kind("slides.pdf", "application/pdf") == MaterialKind.DOCUMENT
