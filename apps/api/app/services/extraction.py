from __future__ import annotations

import subprocess
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile

from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from app.models.entities import Material
from app.models.enums import MaterialKind
from app.providers.base import TranscriptionSegment
from app.providers.factory import get_stt_provider


@dataclass
class ExtractedSection:
    text: str
    page_number: int | None = None
    slide_number: int | None = None
    section_heading: str | None = None
    start_second: int | None = None
    end_second: int | None = None


@dataclass
class ExtractedContent:
    sections: list[ExtractedSection]
    transcript_segments: list[TranscriptionSegment]

    @property
    def full_text(self) -> str:
        return "\n\n".join(section.text for section in self.sections if section.text.strip())


class MaterialExtractor:
    def __init__(self) -> None:
        self.stt_provider = get_stt_provider()

    def extract(self, material: Material, content: bytes) -> ExtractedContent:
        extension = Path(material.file_name).suffix.lower()
        if material.source_kind in {MaterialKind.AUDIO, MaterialKind.VIDEO}:
            return self._extract_transcript(extension, content)
        if extension in {".txt", ".md"}:
            text = content.decode("utf-8", errors="ignore")
            return ExtractedContent(sections=[ExtractedSection(text=text)], transcript_segments=[])
        if extension == ".pdf":
            reader = PdfReader(BytesIO(content))
            sections = [
                ExtractedSection(text=(page.extract_text() or "").strip(), page_number=index + 1)
                for index, page in enumerate(reader.pages)
            ]
            return ExtractedContent(sections=sections, transcript_segments=[])
        if extension == ".docx":
            document = Document(BytesIO(content))
            paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
            return ExtractedContent(
                sections=[ExtractedSection(text="\n".join(paragraphs), section_heading="DOCX body")],
                transcript_segments=[],
            )
        if extension == ".pptx":
            presentation = Presentation(BytesIO(content))
            sections: list[ExtractedSection] = []
            for index, slide in enumerate(presentation.slides, start=1):
                texts = [
                    shape.text_frame.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip()
                ]
                sections.append(ExtractedSection(text="\n".join(texts), slide_number=index))
            return ExtractedContent(sections=sections, transcript_segments=[])
        if extension == ".doc":
            return self._extract_with_binary(content, extension, "catdoc")
        if extension == ".ppt":
            return self._extract_with_binary(content, extension, "catppt")
        raise ValueError(f"Unsupported material type: {extension}")

    def _extract_transcript(self, extension: str, content: bytes) -> ExtractedContent:
        with NamedTemporaryFile(delete=True, suffix=extension) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            transcript_segments = self.stt_provider.transcribe(Path(temp_file.name))
        sections = [
            ExtractedSection(
                text=segment.text,
                start_second=segment.start_second,
                end_second=segment.end_second,
            )
            for segment in transcript_segments
        ]
        return ExtractedContent(sections=sections, transcript_segments=transcript_segments)

    def _extract_with_binary(self, content: bytes, extension: str, command: str) -> ExtractedContent:
        with NamedTemporaryFile(delete=True, suffix=extension) as temp_file:
            temp_file.write(content)
            temp_file.flush()
            try:
                result = subprocess.run([command, temp_file.name], capture_output=True, text=True, check=False)
            except FileNotFoundError as exc:
                raise ValueError(
                    f"{command} is not installed in the API container. Install the legacy extractor package."
                ) from exc
            if result.returncode != 0:
                raise ValueError(result.stderr.strip() or f"{command} failed to process {extension}.")
            return ExtractedContent(
                sections=[ExtractedSection(text=result.stdout.strip(), section_heading=f"{extension} document")],
                transcript_segments=[],
            )


def detect_material_kind(file_name: str, mime_type: str) -> MaterialKind:
    extension = Path(file_name).suffix.lower()
    if extension in {".mp3", ".wav", ".m4a", ".aac"} or mime_type.startswith("audio/"):
        return MaterialKind.AUDIO
    if extension in {".mp4", ".mov", ".mkv"} or mime_type.startswith("video/"):
        return MaterialKind.VIDEO
    return MaterialKind.DOCUMENT
